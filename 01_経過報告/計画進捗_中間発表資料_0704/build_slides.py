"""
進捗発表資料（Word・PowerPoint）生成スクリプト
対象：2026/07/04 卒業研究進捗発表会
"""
import os, copy
from docx import Document
from docx.shared import Pt, Cm, RGBColor, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches, Pt as PPt, Emu as PEmu
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(BASE, "計画進捗・中間発表会・最終発表会用予稿テンプレート.docx")
OUT_WORD = os.path.join(BASE, "計画進捗・中間発表_23F302020.docx")
OUT_PPTX = os.path.join(BASE, "計画進捗・中間発表_23F302020.pptx")

# ─────────────────────────────────────────────────────────────────────────────
# 共通ユーティリティ
# ─────────────────────────────────────────────────────────────────────────────

def set_cell_bg(cell, hex_color):
    """表セルの背景色を設定"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), hex_color)
    tcPr.append(shd)

def set_cell_borders(table):
    """テーブル全体に細いボーダーを設定"""
    tbl = table._tbl
    tblPr = tbl.tblPr
    tblBorders = OxmlElement('w:tblBorders')
    for border_name in ('top', 'left', 'bottom', 'right', 'insideH', 'insideV'):
        border = OxmlElement(f'w:{border_name}')
        border.set(qn('w:val'), 'single')
        border.set(qn('w:sz'), '4')
        border.set(qn('w:space'), '0')
        border.set(qn('w:color'), '2563EB')
        tblBorders.append(border)
    tblPr.append(tblBorders)

# ─────────────────────────────────────────────────────────────────────────────
# Word 予稿生成
# ─────────────────────────────────────────────────────────────────────────────

def build_word():
    doc = Document(TEMPLATE)

    # テンプレートの内容をすべてクリア
    for p in doc.paragraphs:
        p._element.getparent().remove(p._element)
    for t in doc.tables:
        t._element.getparent().remove(t._element)

    def add_para(text, style='Normal', align=None, bold=False, size_pt=None, space_before=None, space_after=None):
        p = doc.add_paragraph(style=style)
        if align:
            p.alignment = align
        if space_before is not None:
            p.paragraph_format.space_before = Pt(space_before)
        if space_after is not None:
            p.paragraph_format.space_after = Pt(space_after)
        run = p.add_run(text)
        run.bold = bold
        if size_pt:
            run.font.size = Pt(size_pt)
        return p

    def add_heading(text, level=1):
        p = doc.add_heading(text, level=level)
        return p

    # ── タイトル ──────────────────────────────────────────────────────────────
    title_p = doc.add_paragraph(style='Title')
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_p.paragraph_format.space_before = Pt(0)
    title_p.paragraph_format.space_after = Pt(2)
    run = title_p.add_run("河川氾濫時における自家用車避難とデマンド交通バス活用の比較シミュレーション")
    run.bold = True
    run.font.size = Pt(13)

    sub_p = add_para("XX研究室　23F302020　氏名：＿＿＿＿＿＿", align=WD_ALIGN_PARAGRAPH.CENTER, size_pt=9, space_before=0, space_after=3)

    # ── 1. 研究背景・目的 ────────────────────────────────────────────────────
    add_heading("研究背景・目的")
    add_para(
        "本研究は，平成27年9月関東・東北豪雨に伴う鬼怒川氾濫を背景とする．"
        "2015年9月10日，茨城県常総市三坂町付近で鬼怒川左岸の堤防が決壊し，市域の広い範囲で浸水被害が発生した．"
        "河川氾濫時には，浸水によって道路が通行困難となるだけでなく，"
        "自家用車避難が集中することで混雑や逃げ遅れが発生する可能性がある．"
        "また，高齢者や自家用車を利用できない住民に対しては，"
        "デマンド交通バスを活用した避難支援の検討が必要である．",
        space_before=0, space_after=2
    )
    add_para(
        "本研究の目的は，浸水時の道路閉鎖を考慮した避難シミュレーションを構築し，"
        "自家用車のみで避難した場合と，デマンド交通バスを活用した場合を比較できる基盤を構築することである．"
        "今回の進捗発表では，Phase 1の静的な避難経路検索と，"
        "Phase 2のSUMO/TraCIによる自家用車避難の交通流シミュレーションについて報告する．",
        space_before=0, space_after=4
    )

    # ── 2. 使用データとシミュレーション構成 ──────────────────────────────────
    add_heading("使用データとシミュレーション構成")
    add_para(
        "本研究では，浸水範囲，道路ネットワーク，行政区域，人口メッシュ，避難施設等を用いる（表1）．"
        "浸水深0.5m以上に相当する範囲と道路リンクが重なる場合，その道路を通行困難として扱う．"
        "Phase 2では実データをそのままSUMOへ入れるのではなく，"
        "研究用に加工した派生データ（net.xml・closure_timeline_sumo.json等）を入力として用いる．"
        "SUMO上の0秒は車両の出発開始であり，浸水時系列t0～t7は0～21,600秒へ線形圧縮した．",
        space_before=0, space_after=3
    )

    # 表1
    cap1 = add_para("表1．使用データとSUMO上での役割", space_before=0, space_after=1)
    cap1.alignment = WD_ALIGN_PARAGRAPH.CENTER
    t1_data = [
        ["データ", "研究内での役割"],
        ["浸水KML・A31a", "時刻別道路閉鎖の作成"],
        ["道路NW（GraphML）", "SUMOネットワーク作成"],
        ["250m人口メッシュ", "出発地40地点と車両数設定"],
        ["避難施設データ", "安全避難所19件の目的地設定"],
        ["閉鎖道路時系列", "TraCIによる動的道路閉鎖"],
    ]
    tbl1 = doc.add_table(rows=len(t1_data), cols=2)
    tbl1.style = 'Table Grid'
    set_cell_borders(tbl1)
    for ri, row in enumerate(t1_data):
        for ci, text in enumerate(row):
            cell = tbl1.cell(ri, ci)
            cell.text = text
            r = cell.paragraphs[0].runs[0]
            r.font.size = Pt(7.5)
            if ri == 0:
                r.bold = True
                set_cell_bg(cell, "2563EB")
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    doc.add_paragraph()

    # ── 3. Phase 1：静的な避難経路検索 ──────────────────────────────────────
    add_heading("Phase 1：静的な避難経路検索")
    add_para(
        "Phase 1では，浸水範囲と道路ネットワークを重ね合わせ，"
        "通行困難となる道路を除外したうえで，出発地から避難所までの経路検索を行った．"
        "これは道路閉鎖後に避難所への経路が存在するかを確認する静的なネットワーク分析である．"
        "常総市を中心事例とし，茨城県内の対象41市区町村へ拡張して"
        "市区町村別の避難経路シナリオHTMLを生成した．"
        "ただしPhase 1では車両同士の混雑や到着時間は扱わない．",
        space_before=0, space_after=4
    )

    # ── 4. Phase 2：SUMOによる自家用車避難シミュレーション ─────────────────
    add_heading("Phase 2：SUMOによる自家用車避難シミュレーション")
    add_para(
        "Phase 2では，Phase 1の道路閉鎖条件をSUMO 1.26.0/TraCIへ接続し，"
        "自家用車のみで避難するシナリオAの交通流シミュレーションを構築した．"
        "出発地・避難所をSUMO edgeへスナップし，"
        "t0〜t7の閉鎖時刻に応じてTraCIでSUMO edgeを通行止めにした．"
        "道路閉鎖後は再経路探索を行い，到着・未到着・混雑を記録した（図1）．"
        "常総市では250m人口メッシュ40地点を出発地とし，smallは40台，10pctは120台，fullは1,001台を投入した．"
        "small・10pctは最初の閉鎖時点t0（SUMO上789秒）より前に全車両が到着するため，"
        "閉鎖後の挙動確認にはfull試行を主に用いる．",
        space_before=0, space_after=3
    )

    # 図1（テキストボックス風の囲み）
    fig_box = doc.add_paragraph()
    fig_box.paragraph_format.space_before = Pt(0)
    fig_box.paragraph_format.space_after  = Pt(1)
    fig_box.paragraph_format.left_indent  = Cm(0.5)
    fig_box.paragraph_format.right_indent = Cm(0.5)
    fig_box.add_run(
        "実データ（浸水・道路・人口・避難所）\n"
        "　　↓　　研究用に加工（派生データ）\n"
        "Phase 1：浸水ポリゴン × 道路エッジ → 閉鎖タイムライン → 静的経路検索\n"
        "　　↓　　SUMOネットワーク変換・edge対応\n"
        "Phase 2：車両投入 → TraCI動的閉鎖 → 再経路探索 → 到着/未到着/混雑記録"
    ).font.size = Pt(8.5)

    cap_fig1 = add_para("図1．Phase 1とPhase 2の処理の流れ", space_before=1, space_after=4)
    cap_fig1.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # ── 5. 結果 ──────────────────────────────────────────────────────────────
    add_heading("結果")
    add_para(
        "常総市シナリオAの実行結果を表2に示す．"
        "small・10pct試行では全車両が到着した．"
        "full試行では1,001台のうち987台が到着し，14台が未到着となった．"
        "未到着14台は閉鎖済みの出発edgeから発車できなかった車両であり，逃げ遅れ候補として扱う．"
        "全41市区町村の10pct試行では逃げ遅れ合計0台であった．",
        space_before=0, space_after=3
    )

    # 表2
    doc.add_paragraph().add_run().add_break(WD_BREAK.COLUMN)
    cap2 = add_para("表2．常総市シナリオAのSUMO実行結果", space_before=0, space_after=1)
    cap2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    t2_data = [
        ["試行", "車両数", "到着/未到着", "到着率"],
        ["small", "40", "40/0", "100%"],
        ["10pct", "120", "120/0", "100%"],
        ["full", "1,001", "987/14", "98.6%"],
    ]
    tbl2 = doc.add_table(rows=len(t2_data), cols=4)
    tbl2.style = 'Table Grid'
    set_cell_borders(tbl2)
    for ri, row in enumerate(t2_data):
        for ci, text in enumerate(row):
            cell = tbl2.cell(ri, ci)
            cell.text = text
            r = cell.paragraphs[0].runs[0]
            r.font.size = Pt(8)
            if ri == 0:
                r.bold = True
                set_cell_bg(cell, "2563EB")
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            if ri > 0 and ci in (1, 2, 3):
                cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_para(
        "smallは動作確認，10pctは縮小比較，fullは常総市主結果として扱う．",
        space_before=2, space_after=3
    )
    doc.add_paragraph()

    # 主要避難路
    add_para(
        "また，full試行における主要避難路別混雑を確認したところ，"
        "県道357号で最大137台走行・停止105台・最低速度1.64 m/sが記録された．"
        "国道294号でも最大61台・停止7台が確認され，自家用車避難集中時に主要路線の速度低下が生じることを確認した．",
        space_before=0, space_after=4
    )

    # ── 6. 考察・今後の予定 ─────────────────────────────────────────────────
    add_heading("考察・今後の予定")
    add_para(
        "Phase 1とPhase 2は同じ道路閉鎖を扱うが，評価内容が異なる．"
        "Phase 1は閉鎖後の経路存在性を調べる静的分析であり，"
        "Phase 2は車両を実際に走らせる動的分析である．"
        "そのためPhase 1の到達不可人口とPhase 2の未到着車両数を直接比較しない．"
        "Phase 2により，自家用車のみで避難した場合の交通流ベースラインを確立した．"
        "今後のPhase 3では，デマンド交通バスを導入したシナリオBを作成し，"
        "到着率・未到着・主要避難路混雑・バス輸送人数をPhase 2と比較する予定である．",
        space_before=0, space_after=4
    )

    # ── 7. 結論 ──────────────────────────────────────────────────────────────
    add_heading("結論")
    add_para(
        "本研究では，2015年鬼怒川氾濫を背景に，浸水時の道路閉鎖を考慮した避難シミュレーションを構築した．"
        "Phase 1では静的な避難経路検索を行い，Phase 2ではSUMO/TraCIを用いた"
        "自家用車避難の交通流シミュレーションを実装した．"
        "常総市full試行では1,001台中987台が到着し，14台が閉鎖済み出発edgeの影響を受けた．"
        "この結果を，今後デマンド交通バスを導入するPhase 3の比較基準として固定する．",
        space_before=0, space_after=4
    )

    # ── 参考文献 ─────────────────────────────────────────────────────────────
    add_heading("参考文献")
    refs = [
        "[1] 国土地理院：平成27年9月関東・東北豪雨の情報，"
        "https://www.gsi.go.jp/BOUSAI/H27.taihuu18gou.html，閲覧日：2026年5月20日．",
        "[2] 国土交通省関東地方整備局下館河川事務所：鬼怒川緊急対策プロジェクト，"
        "https://www.ktr.mlit.go.jp/shimodate/shimodate00166.html，閲覧日：2026年5月20日．",
        "[3] Lopez, P. A. et al.: Microscopic Traffic Simulation using SUMO, "
        "2018 21st International Conference on Intelligent Transportation Systems, pp.2575-2582, 2018.",
        "[4] Alvarez Lopez, P. et al.: Simulation of Urban Mobility (SUMO) (1.26.0), "
        "Zenodo, 2026. https://doi.org/10.5281/zenodo.13907886",
    ]
    for ref in refs:
        p = add_para(ref, space_before=0, space_after=1)
        p.paragraph_format.left_indent = Cm(0.5)
        p.paragraph_format.first_line_indent = Cm(-0.5)

    doc.save(OUT_WORD)
    print(f"[Word] 保存完了: {OUT_WORD}")


# ─────────────────────────────────────────────────────────────────────────────
# PowerPoint スライド生成
# ─────────────────────────────────────────────────────────────────────────────

# カラーパレット
C_NAVY    = PRGBColor(0x1E, 0x40, 0xAF)   # ネイビーブルー（ヘッダー）
C_BLUE    = PRGBColor(0x25, 0x63, 0xEB)   # ブルー（アクセント）
C_LIGHT   = PRGBColor(0xDB, 0xEA, 0xFE)   # 薄青（塗り）
C_WHITE   = PRGBColor(0xFF, 0xFF, 0xFF)
C_DARK    = PRGBColor(0x1E, 0x29, 0x3B)   # 濃いグレー（本文）
C_GRAY    = PRGBColor(0x64, 0x74, 0x8B)   # グレー（サブ）
C_AMBER   = PRGBColor(0xF5, 0x9E, 0x0B)   # アンバー（強調）


def px(cm_val):
    """cm → Emu"""
    return Cm(cm_val)


def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None, line_width=None):
    from pptx.util import Pt as PPt2
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        px(left), px(top), px(width), px(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        if line_width:
            shape.line.width = PPt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                 wrap=True, font_name="游ゴシック"):
    txBox = slide.shapes.add_textbox(px(left), px(top), px(width), px(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PPt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txBox


def add_table_slide(slide, data, left, top, width, height,
                    col_widths=None, header_color=C_NAVY, font_size=14):
    """スライド上に表を追加"""
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, px(left), px(top), px(width), px(height)).table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = px(w)
    for ri, row in enumerate(data):
        for ci, text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = PPt(font_size)
            run.font.name = "游ゴシック"
            if ri == 0:
                run.font.bold = True
                run.font.color.rgb = C_WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_LIGHT
    return table


def make_slide_header(slide, title_text, sub_text=None):
    """スライドのヘッダーバーを追加"""
    # ヘッダーバー背景
    add_rect(slide, 0, 0, 33.87, 2.2, C_NAVY)
    # タイトルテキスト
    add_text_box(slide, title_text, 0.5, 0.25, 32, 1.3,
                 font_size=24, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if sub_text:
        add_text_box(slide, sub_text, 0.5, 1.4, 32, 0.6,
                     font_size=13, color=PRGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.LEFT)
    # 下部アクセントライン
    add_rect(slide, 0, 18.5, 33.87, 0.15, C_BLUE)
    # スライド番号エリア（右下）
    add_text_box(slide, "", 30, 18.0, 3, 0.5, font_size=11, color=C_GRAY, align=PP_ALIGN.RIGHT)


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=15, color=C_DARK, indent="・ "):
    txBox = slide.shapes.add_textbox(px(left), px(top), px(width), px(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = PPt(4)
        run = p.add_run()
        run.text = indent + item
        run.font.size = PPt(font_size)
        run.font.name = "游ゴシック"
        run.font.color.rgb = color


def build_pptx():
    prs = Presentation()
    prs.slide_width  = Cm(33.87)   # 16:9 widescreen
    prs.slide_height = Cm(19.05)

    blank = prs.slide_layouts[6]  # blank

    # ──────────────────────────────────────────────────────────────────────────
    # Slide 1: タイトル
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, 33.87, 19.05, C_NAVY)             # 全面背景
    add_rect(sl, 0, 0, 33.87, 19.05,
             PRGBColor(0x1E, 0x40, 0xAF))                 # グラデーション風の重ね

    # 中央白帯
    add_rect(sl, 1.5, 4.5, 30.87, 8.5, C_WHITE)
    # アクセントライン
    add_rect(sl, 1.5, 4.5, 0.4, 8.5, C_AMBER)

    add_text_box(sl,
        "河川氾濫時における自家用車避難と\nデマンド交通バス活用の比較シミュレーション",
        2.5, 5.0, 29, 4.5,
        font_size=28, bold=True, color=C_DARK, align=PP_ALIGN.LEFT)

    add_text_box(sl, "卒業研究 進捗発表会", 2.5, 9.5, 29, 0.8,
                 font_size=16, color=C_GRAY, align=PP_ALIGN.LEFT)
    add_text_box(sl, "XX研究室　氏名：＿＿＿＿＿＿", 2.5, 10.3, 29, 0.8,
                 font_size=16, color=C_DARK, align=PP_ALIGN.LEFT)
    add_text_box(sl, "2026年07月04日", 2.5, 11.2, 29, 0.8,
                 font_size=14, color=C_GRAY, align=PP_ALIGN.LEFT)

    # ──────────────────────────────────────────────────────────────────────────
    # Slide 2: 研究背景
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, 33.87, 19.05, PRGBColor(0xF8, 0xFA, 0xFC))
    make_slide_header(sl, "研究背景", "2015年 鬼怒川氾濫を背景とした避難シミュレーション研究")

    # キーワードカード
    cards = [
        ("堤防決壊", "2015年9月10日\n茨城県常総市三坂町付近\n鬼怒川左岸が決壊"),
        ("浸水・道路閉鎖", "広範囲の浸水により\n道路が通行困難になる"),
        ("避難集中", "自家用車避難が集中し\n混雑・逃げ遅れの恐れ"),
        ("要支援者", "高齢者・車非保有者への\n避難支援が必要"),
    ]
    for i, (title, body) in enumerate(cards):
        x = 1.0 + i * 7.9
        add_rect(sl, x, 2.8, 7.5, 7.0, C_WHITE, C_BLUE, 1.5)
        add_rect(sl, x, 2.8, 7.5, 1.4, C_BLUE)
        add_text_box(sl, title, x+0.2, 2.9, 7.1, 1.2,
                     font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, body, x+0.2, 4.4, 7.1, 4.5,
                     font_size=14, color=C_DARK, align=PP_ALIGN.CENTER)

    add_text_box(sl,
        "→ 避難シミュレーションによる定量的な評価が必要",
        2.0, 10.5, 30, 1.0,
        font_size=17, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────────────────
    # Slide 3: 研究目的
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, 33.87, 19.05, PRGBColor(0xF8, 0xFA, 0xFC))
    make_slide_header(sl, "研究目的")

    add_rect(sl, 1.5, 2.5, 30.87, 3.0, C_LIGHT, C_BLUE, 1)
    add_text_box(sl,
        "浸水時の道路閉鎖を考慮した避難シミュレーションを構築し、\n"
        "自家用車避難とデマンド交通バス活用を定量的に比較する",
        2.0, 2.7, 30, 2.5,
        font_size=19, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    add_text_box(sl, "今回の発表範囲", 1.5, 6.0, 30, 0.8,
                 font_size=16, bold=True, color=C_NAVY)

    phases = [
        ("Phase 1", "静的な避難経路検索\n（浸水×道路ネットワーク → 経路の有無を確認）", C_BLUE),
        ("Phase 2", "SUMOによる自家用車避難シミュレーション\n（車両走行・動的道路閉鎖・混雑評価）", C_NAVY),
        ("Phase 3\n（予定）", "デマンド交通バス導入・シナリオA/B比較\n（未着手）", C_GRAY),
    ]
    for i, (ph, desc, col) in enumerate(phases):
        x = 1.5 + i * 10.6
        add_rect(sl, x, 7.0, 10.0, 4.5, C_WHITE, col, 2)
        add_rect(sl, x, 7.0, 10.0, 1.4, col)
        add_text_box(sl, ph, x+0.2, 7.1, 9.6, 1.2,
                     font_size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, desc, x+0.2, 8.6, 9.6, 3.0,
                     font_size=13, color=C_DARK, align=PP_ALIGN.CENTER)

        if i < 2:
            add_text_box(sl, "→", x + 10.1, 8.5, 0.8, 1.0,
                         font_size=20, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────────────────
    # Slide 4: 研究全体の構成
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, 33.87, 19.05, PRGBColor(0xF8, 0xFA, 0xFC))
    make_slide_header(sl, "研究全体の構成")

    flow_items = [
        ("実データ", "浸水KML・A31a\n道路NW・人口\n避難所・自動車保有", C_GRAY, C_WHITE),
        ("Phase 1", "浸水×道路\n→ 閉鎖タイムライン\n→ 静的経路検索", C_BLUE, C_WHITE),
        ("Phase 2", "SUMO変換\n→ TraCI動的閉鎖\n→ 交通流評価", C_NAVY, C_WHITE),
        ("Phase 3\n（予定）", "バス投入\n→ シナリオB\n→ A/B比較", C_AMBER, C_DARK),
    ]
    box_w = 6.8
    for i, (label, desc, bg, fc) in enumerate(flow_items):
        x = 1.0 + i * 8.1
        add_rect(sl, x, 3.0, box_w, 9.0, bg)
        add_text_box(sl, label, x+0.1, 3.2, box_w-0.2, 1.6,
                     font_size=16, bold=True, color=fc, align=PP_ALIGN.CENTER)
        add_text_box(sl, desc, x+0.1, 5.0, box_w-0.2, 6.5,
                     font_size=13, color=fc, align=PP_ALIGN.CENTER)
        if i < 3:
            add_text_box(sl, "→", x + box_w + 0.1, 6.5, 0.9, 1.5,
                         font_size=22, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)

    add_text_box(sl, "完了", 2.2, 12.5, 6.0, 0.8,
                 font_size=13, bold=True, color=PRGBColor(0x16, 0xA3, 0x4A), align=PP_ALIGN.CENTER)
    add_text_box(sl, "完了", 10.3, 12.5, 6.0, 0.8,
                 font_size=13, bold=True, color=PRGBColor(0x16, 0xA3, 0x4A), align=PP_ALIGN.CENTER)
    add_text_box(sl, "次フェーズ", 18.4, 12.5, 6.0, 0.8,
                 font_size=13, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────────────────
    # Slide 5: 使用データ
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, 33.87, 19.05, PRGBColor(0xF8, 0xFA, 0xFC))
    make_slide_header(sl, "使用データ", "実データ → 派生データ → SUMO入力 の流れ")

    tbl_data = [
        ["用途", "使用データ", "SUMO上での役割"],
        ["浸水範囲", "GSI KML・A31a浸水想定区域", "時刻別道路閉鎖の作成"],
        ["道路", "道路ネットワーク（GraphML）", "SUMOネットワーク joso.net.xml に変換"],
        ["対象地域", "N03 行政区域", "市区町村の抽出"],
        ["出発地", "250m人口メッシュ", "避難車両の発生地点"],
        ["車両数", "人口・自動車保有台数", "small / 10pct / full の車両数換算"],
        ["目的地", "避難施設データ", "安全避難所のSUMO edgeへのスナップ"],
        ["道路閉鎖", "Phase 1 閉鎖タイムライン", "TraCIによる動的通行止め"],
    ]
    add_table_slide(sl, tbl_data, 1.0, 2.5, 31.87, 15.0,
                    col_widths=[4.5, 12.0, 15.37], font_size=14)

    # ──────────────────────────────────────────────────────────────────────────
    # Slide 6: データからSUMO入力まで
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, 33.87, 19.05, PRGBColor(0xF8, 0xFA, 0xFC))
    make_slide_header(sl, "実データからSUMO入力まで",
                      "実データをそのまま動かしているのではなく、研究用に加工した派生データを入力として使用")

    # フロー図（左→右）
    flow = [
        ("実データ", "浸水KML\nA31a GML\n道路NW (GraphML)\n人口メッシュ\n避難施設DB", C_GRAY),
        ("派生データ", "closure_timeline.json\njoso.net.xml\nagent_origins.csv\nshelters_sumo.csv", C_BLUE),
        ("SUMO入力\n（TraCI）", "edge閉鎖\n車両経路XML\n時刻別道路閉鎖", C_NAVY),
        ("シミュレーション\n結果", "到着/未到着\n混雑ログ\n避難完了時刻", C_AMBER),
    ]
    bw, bh = 7.0, 9.5
    for i, (title, body, col) in enumerate(flow):
        x = 0.8 + i * 8.1
        add_rect(sl, x, 2.8, bw, bh, col)
        add_text_box(sl, title, x+0.1, 3.0, bw-0.2, 1.5,
                     font_size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, body, x+0.3, 4.8, bw-0.6, 7.0,
                     font_size=13, color=C_WHITE, align=PP_ALIGN.LEFT)
        if i < 3:
            add_text_box(sl, "→", x + bw + 0.1, 6.5, 0.8, 1.5,
                         font_size=22, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────────────────
    # Slide 7: Phase 1の内容
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, 33.87, 19.05, PRGBColor(0xF8, 0xFA, 0xFC))
    make_slide_header(sl, "Phase 1：静的な避難経路検索")

    # 左：内容
    add_rect(sl, 1.0, 2.5, 14.0, 14.5, C_WHITE, C_BLUE, 1)
    add_text_box(sl, "実施内容", 1.2, 2.7, 13.6, 0.9,
                 font_size=14, bold=True, color=C_NAVY)
    add_bullet_box(sl, [
        "浸水ポリゴン × 道路リンクで閉鎖候補を抽出",
        "浸水深0.5m以上の道路を通行困難として除外",
        "出発地から避難所への最短経路を探索",
        "常総市を中心に茨城県内41市区町村へ拡張",
        "市区町村別シナリオHTMLを生成",
    ], 1.2, 3.7, 13.5, 9.0, font_size=14)

    add_rect(sl, 1.0, 12.8, 14.0, 3.5, PRGBColor(0xFE, 0xF3, 0xC7), C_AMBER, 1)
    add_text_box(sl, "Phase 1の注意点", 1.2, 13.0, 13.5, 0.9,
                 font_size=13, bold=True, color=C_AMBER)
    add_text_box(sl,
        "車両混雑・停止・到着時間は扱わない\n"
        "「経路が存在するか」を確認する静的分析",
        1.2, 13.8, 13.5, 1.8, font_size=12.5, color=C_DARK)

    # 右：フロー図
    add_rect(sl, 16.5, 2.5, 16.3, 14.5, C_WHITE, C_BLUE, 1)
    add_text_box(sl, "処理フロー", 16.7, 2.7, 15.9, 0.9,
                 font_size=14, bold=True, color=C_NAVY)

    steps = [
        ("浸水ポリゴン + 道路エッジ", C_BLUE),
        ("↓ 空間結合（GIS）", C_GRAY),
        ("道路閉鎖タイムライン\n（t0〜t7）", C_NAVY),
        ("↓ Dijkstra 経路探索", C_GRAY),
        ("経路あり / 経路なし\nを判定 → HTML生成", C_NAVY),
    ]
    y = 3.9
    for step_text, col in steps:
        if step_text.startswith("↓"):
            add_text_box(sl, step_text, 17.5, y, 14, 0.6,
                         font_size=13, color=C_GRAY, align=PP_ALIGN.CENTER)
            y += 0.6
        else:
            add_rect(sl, 17.0, y, 15.0, 1.8, col)
            add_text_box(sl, step_text, 17.2, y+0.2, 14.6, 1.4,
                         font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            y += 2.0

    # ──────────────────────────────────────────────────────────────────────────
    # Slide 8: Phase 2の方法
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, 33.87, 19.05, PRGBColor(0xF8, 0xFA, 0xFC))
    make_slide_header(sl, "Phase 2：SUMO交通流シミュレーション",
                      "SUMO 1.26.0 + TraCI による動的道路閉鎖シミュレーション")

    steps8 = [
        ("① ネットワーク変換",   "GraphML → OSM XML → joso.net.xml\n（通常edge 49,356件）"),
        ("② edge ID 対応",       "Phase1 閉鎖edge 764件 → SUMO edge（全matched）\nSUMO segment 計3,158件"),
        ("③ 出発地・避難所スナップ","250mメッシュ重心 → 最寄りSUMO edge\n避難所 → 安全避難所19件をスナップ"),
        ("④ TraCI動的道路閉鎖",  "SUMO 0秒＝出発開始\nt0〜t7を789〜21,600秒に対応"),
        ("⑤ 再経路探索・記録",    "閉鎖後にSUMOが代替経路を探索\n到着/未到着/混雑を60秒ごとに記録"),
    ]
    step_w, step_h = 30.87, 2.4
    for i, (title, desc) in enumerate(steps8):
        y = 2.5 + i * 2.8
        add_rect(sl, 1.0, y, 0.6, step_h, C_NAVY)
        add_rect(sl, 1.8, y, step_w, step_h, C_WHITE, C_BLUE, 1)
        add_text_box(sl, title, 2.0, y+0.1, 10, 0.9,
                     font_size=15, bold=True, color=C_NAVY)
        add_text_box(sl, desc, 2.0, y+1.0, step_w-0.4, 1.2,
                     font_size=13, color=C_DARK)

    # ──────────────────────────────────────────────────────────────────────────
    # Slide 9: Phase 2の前提条件
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, 33.87, 19.05, PRGBColor(0xF8, 0xFA, 0xFC))
    make_slide_header(sl, "Phase 2 前提条件：時間軸・出発地・台数",
                      "結果の解釈で質問されやすい設定を明示")

    time_data = [
        ["項目", "設定", "解釈上の注意"],
        ["シミュレーション時間", "0〜21,600秒（6時間）", "実時間を6時間に線形圧縮"],
        ["SUMO 0秒", "車両の出発開始", "浸水時系列のt0ではない"],
        ["最初の閉鎖t0", "SUMO上 789秒", "small/10pctはt0前に全車到着"],
        ["最終時点t7", "SUMO上 21,600秒", "閉鎖後挙動はfull試行で確認"],
    ]
    add_table_slide(sl, time_data, 1.0, 2.5, 15.5, 8.0,
                    col_widths=[4.2, 5.4, 5.9], font_size=12.5)

    origin_data = [
        ["項目", "設定", "根拠・作り方"],
        ["出発地", "40地点", "250m人口メッシュ代表点"],
        ["目的地", "安全避難所19件", "洪水対応避難所をSUMO edgeへスナップ"],
        ["small", "40台", "各メッシュ1台"],
        ["10pct", "120台", "fullを各メッシュごとに約1割へ縮小・切り上げ"],
        ["full", "1,001台", "人口÷2.3人/世帯をメッシュごとに切り上げ"],
    ]
    add_table_slide(sl, origin_data, 17.4, 2.5, 15.5, 9.6,
                    col_widths=[3.2, 3.8, 8.5], font_size=12)

    add_rect(sl, 1.0, 13.0, 31.87, 3.8, PRGBColor(0xFE, 0xF3, 0xC7), C_AMBER, 1.2)
    add_text_box(sl, "可視化の注意", 1.5, 13.2, 30.8, 0.8,
                 font_size=15, bold=True, color=C_AMBER)
    add_text_box(sl,
        "現在のHTMLアニメーションはsmall/10pctの車両移動確認が中心であり、両試行はt0前に到着する。\n"
        "道路閉鎖後の交通挙動は、full試行の結果表・混雑指標・閉鎖時系列を併用して説明する。",
        1.5, 14.1, 30.8, 1.9, font_size=13.5, color=C_DARK)

    # ──────────────────────────────────────────────────────────────────────────
    # Slide 10: 常総市シナリオAの結果
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, 33.87, 19.05, PRGBColor(0xF8, 0xFA, 0xFC))
    make_slide_header(sl, "Phase 2 結果：常総市シナリオA（自家用車のみ）")

    result_data = [
        ["試行",    "車両数",   "到着台数",  "未到着台数",  "到着率",  "位置づけ"],
        ["small",  "40",     "40",       "0",          "100%",   "動作確認"],
        ["10pct",  "120",    "120",      "0",          "100%",   "縮小比較"],
        ["full",   "1,001",  "987",      "14",         "98.6%",  "常総市主結果"],
    ]
    add_table_slide(sl, result_data, 1.5, 2.5, 30.87, 7.5,
                    col_widths=[4.5, 4.5, 4.5, 4.5, 4.0, 8.87], font_size=15)

    # 未到着説明ボックス
    add_rect(sl, 1.5, 10.5, 30.87, 2.5, PRGBColor(0xFE, 0xF3, 0xC7), C_AMBER, 1.5)
    add_text_box(sl,
        "未到着14台について",
        2.0, 10.7, 29, 0.8, font_size=15, bold=True, color=C_AMBER)
    add_text_box(sl,
        "閉鎖済みの出発edgeから発車できなかった車両として記録。\n"
        "プログラムエラーではなく、今回の閉鎖条件下での「逃げ遅れ候補」として扱う。",
        2.0, 11.5, 29, 1.3, font_size=14, color=C_DARK)

    # 全域結果
    add_rect(sl, 1.5, 13.5, 30.87, 3.8, C_LIGHT, C_BLUE, 1)
    add_text_box(sl, "全41市区町村（10pct試行）の結果",
                 2.0, 13.7, 29, 0.9, font_size=15, bold=True, color=C_NAVY)
    add_bullet_box(sl, [
        "全41市区町村でsmall・10pct完了",
        "10pct試行：計23,054台、逃げ遅れ合計 0台",
        "full試行：代表6市区町村で実行（全員到着）",
    ], 2.0, 14.6, 29, 2.5, font_size=14, indent="● ")

    # ──────────────────────────────────────────────────────────────────────────
    # Slide 11: 混雑・主要避難路分析
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, 33.87, 19.05, PRGBColor(0xF8, 0xFA, 0xFC))
    make_slide_header(sl, "Phase 2 結果：主要避難路別の混雑（full試行）")

    cong_data = [
        ["路線",          "最大走行台数", "最大停止台数", "最低平均速度", "約 km/h"],
        ["県道357号\n（谷和原筑西線）", "137台", "105台", "1.64 m/s", "約 5.9 km/h"],
        ["国道294号",    "61台",      "7台",       "4.41 m/s",  "約 15.9 km/h"],
    ]
    add_table_slide(sl, cong_data, 1.5, 2.5, 30.87, 7.5,
                    col_widths=[7.5, 5.0, 5.0, 6.0, 7.37], font_size=15)

    add_rect(sl, 1.5, 10.5, 30.87, 2.0, PRGBColor(0xFE, 0xE2, 0xE2), PRGBColor(0xEF, 0x44, 0x44), 1.5)
    add_text_box(sl, "県道357号：停止台数105台、最低速度1.64 m/s（約6 km/h）\n  → 自家用車避難集中時に深刻な渋滞が発生",
                 2.0, 10.7, 29, 1.6, font_size=14, bold=False, color=PRGBColor(0x7F, 0x1D, 0x1D))

    add_rect(sl, 1.5, 13.0, 14.5, 4.5, C_WHITE, C_NAVY, 1)
    add_text_box(sl, "Phase 2で分かること", 2.0, 13.2, 14, 0.8,
                 font_size=14, bold=True, color=C_NAVY)
    add_bullet_box(sl, [
        "到着できたかの判定（到着/未到着）",
        "どの道路に負荷が集中するか",
        "速度低下・停止が発生する区間",
    ], 2.0, 14.0, 13.5, 3.0, font_size=13)

    add_rect(sl, 17.5, 13.0, 14.87, 4.5, C_WHITE, C_AMBER, 1)
    add_text_box(sl, "Phase 2でわからないこと", 18.0, 13.2, 14, 0.8,
                 font_size=14, bold=True, color=C_AMBER)
    add_bullet_box(sl, [
        "車なし世帯・高齢者の避難可否",
        "バスが渋滞を緩和するか（Phase 3で比較）",
        "出発時刻ばらつきの影響",
    ], 18.0, 14.0, 13.5, 3.0, font_size=13)

    # ──────────────────────────────────────────────────────────────────────────
    # Slide 12: Phase 1とPhase 2の違い
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, 33.87, 19.05, PRGBColor(0xF8, 0xFA, 0xFC))
    make_slide_header(sl, "Phase 1 と Phase 2 の違い")

    diff_data = [
        ["観点",          "Phase 1（静的）",              "Phase 2（動的）"],
        ["分析手法",      "道路閉鎖後の最短経路の有無を判定", "車両を実際に走らせて結果を観測"],
        ["渋滞の考慮",    "なし",                         "あり（速度低下・停止が記録される）"],
        ["指標の単位",    "「経路なし」メッシュ数・人口数",   "「実際に未到着」の車両台数"],
        ["常総市の結果",  "到達不可メッシュ17・人口952名",   "full試行で14台が発車できず逃げ遅れ"],
        ["わかること",    "「そもそも道がない」場所",        "「道はあるが渋滞・閉鎖で間に合わない」状況"],
    ]
    add_table_slide(sl, diff_data, 1.0, 2.5, 31.87, 14.0,
                    col_widths=[5.5, 13.0, 13.37], font_size=14)

    add_rect(sl, 1.0, 17.2, 31.87, 1.2, PRGBColor(0xFE, 0xF3, 0xC7), C_AMBER, 1)
    add_text_box(sl,
        "注意：Phase 1の到達不可人数（952名）とPhase 2の未到着車両数（14台）は単位が異なる → 直接比較しない",
        1.5, 17.3, 31, 1.0, font_size=13, bold=True, color=C_AMBER)

    # ──────────────────────────────────────────────────────────────────────────
    # Slide 13: 今後の予定
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, 33.87, 19.05, PRGBColor(0xF8, 0xFA, 0xFC))
    make_slide_header(sl, "今後の予定（Phase 3）",
                      "デマンド交通バスを追加し、自家用車のみのシナリオAと比較する")

    # 比較表
    comp_data = [
        ["比較軸",     "Phase 2（シナリオA）",     "Phase 3（シナリオB）"],
        ["道路上の車両", "自家用車のみ",            "自家用車 ＋ バス"],
        ["対象者",     "車保有世帯のみ追跡",         "車なし・高齢者もバスで追跡"],
        ["主要路渋滞",  "県道357号等で深刻な渋滞を確認", "バス導入による緩和効果を測定"],
        ["結論",       "比較基準（ベースライン）",   "RQへの定量的な回答"],
    ]
    add_table_slide(sl, comp_data, 1.0, 2.5, 31.87, 9.0,
                    col_widths=[5.5, 13.0, 13.37], font_size=14)

    add_text_box(sl, "Phase 3 実装ステップ", 1.0, 12.0, 31.87, 0.9,
                 font_size=15, bold=True, color=C_NAVY)
    steps12 = [
        ("① バス設定仕様固定", "台数・定員・運行方式\n利用者割り当てルール"),
        ("② シナリオB実装", "常総市 small\nSUMO/TraCI バス投入"),
        ("③ A/B 比較", "到着率・混雑・\nバス輸送人数を算出"),
        ("④ 卒論執筆", "Phase 2/3 本文\n第3〜5章"),
    ]
    for i, (step, desc) in enumerate(steps12):
        x = 1.0 + i * 8.1
        add_rect(sl, x, 13.0, 7.5, 5.0, C_NAVY)
        add_text_box(sl, step, x+0.2, 13.2, 7.1, 1.0,
                     font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, desc, x+0.2, 14.5, 7.1, 3.0,
                     font_size=13, color=PRGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)
        if i < 3:
            add_text_box(sl, "→", x + 7.6, 14.8, 0.7, 1.0,
                         font_size=18, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)

    prs.save(OUT_PPTX)
    print(f"[PPTX] 保存完了: {OUT_PPTX}")


# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    build_word()
    build_pptx()
    print("✅ 全ファイル生成完了")
